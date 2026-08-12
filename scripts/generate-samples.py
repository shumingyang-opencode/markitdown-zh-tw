#!/usr/bin/env python3
"""Generate sample files for markitdown-zh-tw conversion case studies.

Creates small, self-contained input files under `worked/conversions/<group>/`
covering: DOCX, XLSX, PPTX, PDF, HTML, CSV, JSON, JPEG (with EXIF), WAV,
ZIP and EPUB. The files are generated locally with open-source Python
libraries (python-docx / openpyxl / python-pptx / reportlab / Pillow / wave).
"""

from __future__ import annotations

import csv
import json
import math
import struct
import wave
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw
import piexif  # type: ignore

try:
    import docx  # python-docx
    from docx.shared import Pt  # type: ignore
except ImportError:  # pragma: no cover
    docx = None

try:
    import openpyxl  # type: ignore
except ImportError:  # pragma: no cover
    openpyxl = None

try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt as PptPt  # type: ignore
except ImportError:  # pragma: no cover
    Presentation = None

try:
    from reportlab.lib.pagesizes import A4  # type: ignore
    from reportlab.platypus import (  # type: ignore
        SimpleDocTemplate,
        Paragraph,
        Spacer,
        Table,
        TableStyle,
    )
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle  # type: ignore
    from reportlab.lib import colors  # type: ignore
    from reportlab.pdfbase import pdfmetrics  # type: ignore
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont  # type: ignore
except ImportError:  # pragma: no cover
    SimpleDocTemplate = None

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "worked" / "conversions"


def w(*parts: str) -> Path:
    """Resolve a path under worked/conversions and ensure its parent exists."""
    p = OUT.joinpath(*parts)
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


# --------------------------------------------------------------------------
# DOCX (Word)
# --------------------------------------------------------------------------
def gen_docx() -> None:
    if docx is None:  # pragma: no cover
        return
    d = docx.Document()
    d.add_heading("MarkItDown 測試文件", level=0)
    d.add_paragraph(
        "這是一份用來示範 MarkItDown 轉換的 Word 文件。"
        "包含標題、段落、項目符號與表格，轉換後應保留結構。"
    )
    d.add_heading("第一個章節", level=1)
    d.add_paragraph("段落文字：markitdown 會保留 heading、list 與 table 等結構。")
    d.add_heading("功能清單", level=2)
    for item in ["PDF 轉換", "Word 轉換", "Excel 轉換", "簡報轉換", "圖片/音訊"]:
        d.add_paragraph(item, style="List Bullet")
    d.add_heading("版本比較表", level=2)
    table = d.add_table(rows=3, cols=3)
    table.style = "Light Grid Accent 1"
    data = [["功能", "markitdown", "textract"], ["PDF", "✅", "✅"], ["DOCX", "✅", "✅"]]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    d.save(str(w("office", "sample-report.docx")))


# --------------------------------------------------------------------------
# XLSX (Excel)
# --------------------------------------------------------------------------
def gen_xlsx() -> None:
    if openpyxl is None:  # pragma: no cover
        return
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Region", "Product", "Units", "Revenue"])
    rows = [
        ["North", "Widget", 120, 2400],
        ["North", "Gadget", 80, 3200],
        ["South", "Widget", 150, 3000],
        ["South", "Gadget", 60, 2400],
        ["East", "Widget", 90, 1800],
    ]
    for r in rows:
        ws.append(r)
    ws.append(["Total", "", 500, 12800])
    ws["A8"] = "註：此為樣本資料，用於示範 Excel → Markdown 的表格保留。"
    wb.save(str(w("office", "sales-data.xlsx")))


# --------------------------------------------------------------------------
# PPTX (PowerPoint)
# --------------------------------------------------------------------------
def gen_pptx() -> None:
    if Presentation is None:  # pragma: no cover
        return
    prs = Presentation()
    # Slide 1: 標題
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "MarkItDown 簡報範例"
    s.placeholders[1].text = "示範 PowerPoint 轉 Markdown\n作者：markitdown-zh-tw"
    # Slide 2: 項目符號
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "為什麼用 Markdown？"
    tf = s.placeholders[1].text_frame
    for i, line in enumerate(
        ["對 LLM 友善", "token 效率高", "保留文件結構"]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
    # Slide 3: 表格
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "支援格式一覽"
    rows, cols = 4, 2
    table = s.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(2)).table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = ["格式", "狀態"][c] if r == 0 else [["PDF", "✅"], ["DOCX", "✅"], ["XLSX", "✅"]][r - 1][c]
    prs.save(str(w("office", "pitch-deck.pptx")))


# --------------------------------------------------------------------------
# PDF
# --------------------------------------------------------------------------
def gen_pdf() -> None:
    if SimpleDocTemplate is None:  # pragma: no cover
        return
    path = str(w("pdf", "report.pdf"))
    # 註冊內建 CID 中文字型，確保 PDF 內嵌正確的 CJK 字形
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    except Exception:  # pragma: no cover
        pass
    doc = SimpleDocTemplate(path, pagesize=A4)
    styles = getSampleStyleSheet()
    title_zh = ParagraphStyle(
        "TitleZh", parent=styles["Title"], fontName="STSong-Light"
    )
    body_zh = ParagraphStyle(
        "BodyZh", parent=styles["BodyText"], fontName="STSong-Light"
    )
    head_zh = ParagraphStyle(
        "HeadZh", parent=styles["Heading2"], fontName="STSong-Light"
    )
    story = [
        Paragraph("MarkItDown 測試報告", title_zh),
        Paragraph("這是一份用來示範 PDF 轉換的報告，包含標題、段落與表格。", body_zh),
        Spacer(1, 12),
        Paragraph("效能數據", head_zh),
        Table(
            [["方法", "耗時(ms)"], ["pdfminer", "120"], ["pdfplumber", "180"]],
            style=TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#22d3ee")),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("FONTNAME", (0, 0), (-1, -1), "STSong-Light"),
                ]
            ),
        ),
    ]
    doc.build(story)


# --------------------------------------------------------------------------
# HTML / CSV / JSON (手寫文字檔)
# --------------------------------------------------------------------------
def gen_text_files() -> None:
    html = """<!DOCTYPE html>
<html lang="zh-Hant">
<head><meta charset="utf-8"><title>MarkItDown 範例網頁</title></head>
<body>
  <h1>MarkItDown 範例網頁</h1>
  <p>這是 <strong>HTML</strong> 範例，用於示範 <code>HtmlConverter</code> 的轉換。</p>
  <h2>連結</h2>
  <ul>
    <li><a href="https://github.com/microsoft/markitdown">MarkItDown GitHub</a></li>
    <li><a href="https://pypi.org/project/markitdown/">PyPI</a></li>
  </ul>
  <h2>表格</h2>
  <table border="1">
    <tr><th>功能</th><th>支援</th></tr>
    <tr><td>標題</td><td>✅</td></tr>
    <tr><td>表格</td><td>✅</td></tr>
  </table>
</body>
</html>"""
    w("web", "example.html").write_text(html, encoding="utf-8")

    with w("web", "data.csv").open("w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        writer.writerows(
            [
                ["name", "role", "city"],
                ["Alice", "Engineer", "Taipei"],
                ["Bob", "Designer", "Hsinchu"],
                ["Carol", "Manager", "Taichung"],
            ]
        )

    w("web", "data.json").write_text(
        json.dumps(
            {"project": "markitdown-zh-tw", "versions": [0, 1, 7], "tags": ["markdown", "llm"]},
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )


# --------------------------------------------------------------------------
# JPEG（含 EXIF）
# --------------------------------------------------------------------------
def gen_image() -> None:
    img = Image.new("RGB", (640, 360), (13, 13, 19))
    d = ImageDraw.Draw(img)
    d.rectangle([20, 20, 620, 340], outline=(34, 211, 238), width=3)
    d.text((40, 40), "MarkItDown 圖片範例", fill=(236, 234, 244))
    d.text((40, 80), "ImageConverter: EXIF + LLM 描述", fill=(163, 161, 182))
    path = str(w("media", "photo.jpg"))
    img.save(path, quality=92)
    exif_bytes = piexif.dump(
        {
            "0th": {
                piexif.ImageIFD.Make: "markitdown-zh-tw",
                piexif.ImageIFD.Model: "sample",
            },
            "Exif": {
                piexif.ExifIFD.DateTimeOriginal: "2026:08:12 10:00:00",
            },
        }
    )
    piexif.insert(exif_bytes, path)


# --------------------------------------------------------------------------
# WAV（音訊）
# --------------------------------------------------------------------------
def gen_wav() -> None:
    path = str(w("media", "tone.wav"))
    sr, duration, freq = 22050, 2, 440
    frames = b"".join(
        struct.pack("<h", int(16000 * math.sin(2 * math.pi * freq * t / sr)))
        for t in range(sr * duration)
    )
    with wave.open(path, "wb") as f:
        f.setnchannels(1)
        f.setsampwidth(2)
        f.setframerate(sr)
        f.writeframes(frames)


# --------------------------------------------------------------------------
# ZIP
# --------------------------------------------------------------------------
def gen_zip() -> None:
    with zipfile.ZipFile(str(w("archive", "archive.zip")), "w") as z:
        z.writestr("readme.txt", "Hello from a ZIP file!\n")
        z.writestr("data/notes.md", "# Notes\n\nZIP 內含多個檔案，會被逐一轉換。\n")
        z.writestr("data/todo.txt", "1. Convert zip\n2. Profit\n")


# --------------------------------------------------------------------------
# EPUB（最小結構）
# --------------------------------------------------------------------------
def gen_epub() -> None:
    container = """<?xml version="1.0"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles>
    <rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>
  </rootfiles>
</container>"""
    opf = """<?xml version="1.0" encoding="utf-8"?>
<package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="bookid">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:title>MarkItDown 迷你電子書</dc:title>
    <dc:creator>markitdown-zh-tw</dc:creator>
    <dc:language>zh-Hant</dc:language>
    <dc:identifier id="bookid">urn:uuid:markitdown-zh-tw-sample</dc:identifier>
  </metadata>
  <manifest>
    <item id="c1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>
  </manifest>
  <spine><itemref idref="c1"/></spine>
</package>"""
    chapter = """<?xml version="1.0" encoding="utf-8"?>
<html xmlns="http://www.w3.org/1999/xhtml">
<head><title>第一章</title></head>
<body>
  <h1>第一章：什麼是 MarkItDown？</h1>
  <p>MarkItDown 是微軟出品的檔案轉 Markdown 工具，適合餵給 LLM 使用。</p>
</body>
</html>"""
    with zipfile.ZipFile(str(w("archive", "book.epub")), "w") as z:
        z.writestr("mimetype", "application/epub+zip", compress_type=zipfile.ZIP_STORED)
        z.writestr("META-INF/container.xml", container)
        z.writestr("OEBPS/content.opf", opf)
        z.writestr("OEBPS/chapter1.xhtml", chapter)


def main() -> None:
    gen_docx()
    gen_xlsx()
    gen_pptx()
    gen_pdf()
    gen_text_files()
    gen_image()
    gen_wav()
    gen_zip()
    gen_epub()
    print(f"樣本檔已產生於 {OUT}")


if __name__ == "__main__":
    main()
